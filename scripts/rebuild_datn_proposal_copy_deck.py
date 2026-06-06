from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs" / "slides" / "DATN_Proposal.pptx"
OUT = ROOT / "docs" / "slides" / "DATN_Proposal - Copy.pptx"
IMG_DIR = ROOT / "reports" / "final_reports" / "src" / "Images"
TEMPLATE_ASSET_DIR = ROOT / "tmp" / "datn_template_assets"
BASE_WIDE_WIDTH = 13.333333
TEMPLATE_WIDTH = 10.0
X_SCALE = TEMPLATE_WIDTH / BASE_WIDE_WIDTH
FONT_SCALE = 0.86


NAVY = RGBColor(16, 42, 67)
BLUE = RGBColor(36, 102, 179)
TEAL = RGBColor(16, 128, 128)
GREEN = RGBColor(85, 128, 44)
GOLD = RGBColor(219, 145, 28)
ORANGE = RGBColor(230, 111, 55)
RED = RGBColor(202, 67, 67)
PURPLE = RGBColor(112, 80, 170)
INK = RGBColor(36, 48, 62)
MUTED = RGBColor(102, 114, 128)
LIGHT_BG = RGBColor(246, 249, 252)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(207, 216, 226)


def _emu(value: float):
    return Inches(value)


def _x(value: float):
    return Inches(value * X_SCALE)


def _y(value: float):
    return Inches(value)


def _w(value: float):
    return Inches(value * X_SCALE)


def _h(value: float):
    return Inches(value)


def _template_asset(name: str) -> Path:
    target = TEMPLATE_ASSET_DIR / name
    if target.exists():
        return target
    TEMPLATE_ASSET_DIR.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(TEMPLATE) as deck:
        source = f"ppt/media/{name}"
        target.write_bytes(deck.read(source))
    return target


def _clear_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _background(slide, prs, *, kind: str = "content"):
    image = "image2.png" if kind == "cover" else "image1.png" if kind == "closing" else "image3.png"
    slide.shapes.add_picture(
        str(_template_asset(image)),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def _add_bg(slide, prs, title: str | None = None, section: str | None = None):
    # Content slides use the original template layout directly:
    # red title band, gold section band, white body, and HUST footer.
    if section:
        box = slide.shapes.add_textbox(_x(0.38), _y(0.87), _w(2.4), _h(0.24))
        p = box.text_frame.paragraphs[0]
        p.text = section.upper()
        p.font.name = "Aptos"
        p.font.size = Pt(8 * FONT_SCALE)
        p.font.bold = True
        p.font.color.rgb = RGBColor(120, 72, 0)
        p.alignment = PP_ALIGN.LEFT
    if title:
        tx = slide.shapes.add_textbox(_x(0.45), _y(0.13), _w(12.35), _h(0.46))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Aptos Display"
        p.font.size = Pt(22 * FONT_SCALE)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

def _text_box(slide, text: str, x: float, y: float, w: float, h: float, *, size=18, color=INK,
              bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Aptos"):
    box = slide.shapes.add_textbox(_x(x), _y(y), _w(w), _h(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size * FONT_SCALE)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullet_box(slide, items: list[str], x: float, y: float, w: float, h: float, *, size=17,
                color=INK, bullet_color=TEAL):
    box = slide.shapes.add_textbox(_x(x), _y(y), _w(w), _h(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size * FONT_SCALE)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def _card(slide, x: float, y: float, w: float, h: float, *, fill=CARD, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, _x(x), _y(y), _w(w), _h(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.1)
    return shp


def _label(slide, text: str, x: float, y: float, w: float, h: float, *, fill=TEAL, color=RGBColor(255, 255, 255),
           size=11):
    shp = _card(slide, x, y, w, h, fill=fill, line=fill)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size * FONT_SCALE)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shp


def _stat_card(slide, label: str, value: str, note: str, x: float, y: float, w: float, h: float, color=TEAL):
    _card(slide, x, y, w, h)
    _text_box(slide, value, x + 0.18, y + 0.16, w - 0.36, 0.45, size=25, color=color, bold=True,
              align=PP_ALIGN.CENTER)
    _text_box(slide, label, x + 0.18, y + 0.72, w - 0.36, 0.28, size=11, color=NAVY, bold=True,
              align=PP_ALIGN.CENTER)
    _text_box(slide, note, x + 0.18, y + 1.04, w - 0.36, h - 1.15, size=9.5, color=MUTED,
              align=PP_ALIGN.CENTER)


def _arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _x(x1), _y(y1), _x(x2), _y(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def _fit_picture(slide, path: Path, x: float, y: float, w: float, h: float):
    pic = slide.shapes.add_picture(str(path), _x(x), _y(y), width=_w(w))
    if pic.height > _h(h):
        ratio = _h(h) / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = _h(h)
    pic.left = _x(x + (w - (pic.width / 914400) / X_SCALE) / 2)
    pic.top = _y(y + (h - pic.height / 914400) / 2)
    return pic


def _bar(slide, label: str, value: float, x: float, y: float, w: float, color, max_value=100.0):
    _text_box(slide, label, x, y, 2.4, 0.28, size=12, color=INK, bold=True)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _x(x + 2.45), _y(y + 0.04), _w(w), _h(0.18))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(225, 232, 240)
    bg.line.fill.background()
    fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _x(x + 2.45), _y(y + 0.04), _w(w * value / max_value), _h(0.18))
    fg.fill.solid()
    fg.fill.fore_color.rgb = color
    fg.line.fill.background()
    _text_box(slide, f"{value:.1f}%", x + 2.55 + w, y - 0.02, 0.75, 0.26, size=11, color=NAVY, bold=True)


def _mini_node(slide, label: str, x: float, y: float, color, *, w: float = 1.16, h: float = 0.44, size: float = 8.6):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _x(x), _y(y), _w(w), _h(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Aptos"
    p.font.size = Pt(size * FONT_SCALE)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    return shp


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    hust_red = RGBColor(196, 0, 0)
    slide.shapes.add_picture(str(_template_asset("image8.png")), _x(0.55), _y(0.42), width=_w(2.35))
    _text_box(slide, "CogMem", 0.72, 1.48, 4.8, 0.7, size=44, color=hust_red, bold=True, font="Aptos Display")
    _text_box(slide, "Cognitive-Grounded Long-Term Memory\nfor Conversational Agents", 0.74, 2.28, 5.4, 0.7,
              size=18, color=hust_red, bold=True)
    _text_box(slide, "Typed memory graph + multi-channel recall\n+ SUM spreading activation", 0.76, 3.32, 5.1, 0.55,
              size=12.5, color=hust_red, bold=True)
    _text_box(slide, "Student: Le Minh Triet", 0.76, 6.34, 4.8, 0.28, size=13.5, color=hust_red, bold=True)
    # Mini pipeline visual, placed in the open center-right of the title template.
    steps = [("Chat", BLUE), ("Graph", TEAL), ("Evidence", GOLD), ("Answer", GREEN)]
    x = 6.65
    for idx, (name, color) in enumerate(steps):
        _label(slide, name, x + idx * 1.22, 4.4, 1.12, 0.38, fill=color, size=8.8)
        if idx < len(steps) - 1:
            _arrow(slide, x + idx * 1.22 + 1.12, 4.59, x + idx * 1.22 + 1.22, 4.59, color=MUTED, width=1.2)
    _text_box(slide, "Defense focus", 6.88, 1.65, 4.8, 0.35, size=16, color=NAVY, bold=True)
    _bullet_box(slide, [
        "What was built: a typed memory graph pipeline.",
        "What was proven: manual LongMemEval and LoCoMo gains.",
        "What remains: temporal errors and habit-specific validation.",
    ], 6.88, 2.08, 4.85, 1.25, size=12.7)
    _text_box(slide, "10-15 minute storyline:\nproblem -> method -> example ->\nevidence -> contribution proof",
              6.88, 5.08, 4.8, 0.78, size=13.5, color=INK, bold=True)
    return slide


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Problem: long conversations break short-context memory", "motivation")
    _text_box(slide, "A useful assistant must remember across many sessions, but the answer is rarely stored as one clean sentence.",
              0.72, 1.42, 11.7, 0.55, size=18, color=INK, bold=True)
    # timeline
    y = 3.0
    _arrow(slide, 0.95, y, 11.85, y, color=MUTED, width=1.4)
    sessions = [
        ("Session 1", "Preference", "likes quiet hotels", ORANGE),
        ("Session 8", "Location", "Seattle trip", TEAL),
        ("Session 17", "Plan", "workshop next month", GOLD),
        ("Session 31", "Outcome", "tool fix worked", GREEN),
        ("Session 47", "Conflict", "updated answer", RED),
    ]
    xs = [1.0, 3.35, 5.6, 7.8, 10.05]
    for x, (s, kind, detail, color) in zip(xs, sessions):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _x(x), _y(y - 0.16), _w(0.32), _h(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        _card(slide, x - 0.33, y + 0.35, 1.58, 1.05)
        _text_box(slide, s, x - 0.2, y + 0.46, 1.3, 0.22, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, kind, x - 0.2, y + 0.72, 1.3, 0.22, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, detail, x - 0.2, y + 0.97, 1.3, 0.28, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    _card(slide, 0.78, 5.0, 3.4, 1.15)
    _card(slide, 4.55, 5.0, 3.4, 1.15)
    _card(slide, 8.32, 5.0, 3.4, 1.15)
    _text_box(slide, "Raw log", 0.98, 5.18, 3.0, 0.25, size=14, color=NAVY, bold=True)
    _text_box(slide, "too long and noisy for every prompt", 0.98, 5.52, 3.0, 0.42, size=12, color=MUTED)
    _text_box(slide, "Vector-only memory", 4.75, 5.18, 3.0, 0.25, size=14, color=NAVY, bold=True)
    _text_box(slide, "semantic match, but weak relation/time control", 4.75, 5.52, 3.0, 0.42, size=12, color=MUTED)
    _text_box(slide, "Automatic judge", 8.52, 5.18, 3.0, 0.25, size=14, color=NAVY, bold=True)
    _text_box(slide, "not reliable enough as final truth", 8.52, 5.52, 3.0, 0.42, size=12, color=MUTED)
    return slide


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Solution overview: Retain -> Graph -> Recall -> Answer", "architecture")
    _fit_picture(slide, IMG_DIR / "cogmem_pipeline_overview.png", 0.55, 1.16, 11.85, 4.45)
    _label(slide, "1. Retain typed facts", 0.78, 5.83, 2.25, 0.42, fill=TEAL)
    _label(slide, "2. Store memory graph", 3.28, 5.83, 2.25, 0.42, fill=GOLD)
    _label(slide, "3. Multi-channel recall", 5.78, 5.83, 2.25, 0.42, fill=BLUE)
    _label(slide, "4. Grounded generation", 8.28, 5.83, 2.45, 0.42, fill=GREEN)
    _text_box(slide, "Feedback loop: each answer can become new retained evidence later.", 2.15, 6.45, 8.8, 0.35,
              size=14, color=INK, align=PP_ALIGN.CENTER)
    return slide


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "What is stored: compressed fact + lossless snippet", "memory representation")
    _card(slide, 0.78, 1.25, 5.3, 4.95, fill=RGBColor(255, 255, 255), line=RGBColor(190, 219, 219))
    _card(slide, 6.65, 1.25, 5.3, 4.95, fill=RGBColor(255, 255, 255), line=RGBColor(232, 203, 150))
    _label(slide, "Layer 1: Narrative fact", 1.08, 1.52, 2.35, 0.42, fill=TEAL)
    _label(slide, "Layer 2: Raw snippet", 6.95, 1.52, 2.1, 0.42, fill=GOLD)
    _text_box(slide, "Structured memory unit", 1.08, 2.18, 4.7, 0.35, size=17, color=NAVY, bold=True)
    _bullet_box(slide, [
        "typed fact: experience / intention / action-effect / ...",
        "entities, dates, confidence and metadata",
        "embedding for semantic retrieval",
        "graph links for multi-hop recall",
    ], 1.08, 2.68, 4.55, 2.1, size=13)
    _text_box(slide, "Example fact", 1.08, 5.05, 4.7, 0.25, size=12, color=TEAL, bold=True)
    _text_box(slide, "User was in Seattle before traveling to Chicago.", 1.08, 5.35, 4.7, 0.42, size=13.5, color=INK)
    _text_box(slide, "Source-grounded detail", 6.95, 2.18, 4.7, 0.35, size=17, color=NAVY, bold=True)
    _bullet_box(slide, [
        "keeps original wording and context",
        "helps generation avoid lossy summaries",
        "supports evidence guard and manual review",
        "useful when dates or durations are implicit",
    ], 6.95, 2.68, 4.55, 2.1, size=13)
    _text_box(slide, "Example snippet", 6.95, 5.05, 4.7, 0.25, size=12, color=GOLD, bold=True)
    _text_box(slide, "\"I stayed in Seattle before flying to Chicago.\"", 6.95, 5.35, 4.7, 0.42, size=13.5, color=INK)
    _arrow(slide, 6.08, 3.73, 6.65, 3.73, color=MUTED, width=1.8)
    _text_box(slide, "Design goal: compact enough to retrieve, faithful enough to answer.", 1.1, 6.47, 10.8, 0.35,
              size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Six typed memory networks", "memory representation")
    _fit_picture(slide, IMG_DIR / "cogmem_memory_graph.png", 0.32, 1.12, 7.0, 4.95)
    cards = [
        ("World", "stable facts", BLUE),
        ("Experience", "dated events", TEAL),
        ("Opinion", "preferences", ORANGE),
        ("Habit", "routines", GREEN),
        ("Intention", "future plans", GOLD),
        ("Action-effect", "precondition -> action -> outcome", RED),
    ]
    for i, (name, desc, color) in enumerate(cards):
        x = 7.55 + (i % 2) * 2.45
        y = 1.25 + (i // 2) * 1.12
        _card(slide, x, y, 2.15, 0.84, fill=RGBColor(255, 255, 255), line=color)
        _text_box(slide, name, x + 0.15, y + 0.12, 1.85, 0.23, size=12.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, desc, x + 0.15, y + 0.42, 1.85, 0.25, size=8.7, color=MUTED, align=PP_ALIGN.CENTER)
    _card(slide, 7.55, 4.92, 4.6, 1.15, fill=RGBColor(255, 252, 244), line=GOLD)
    _text_box(slide, "Honest scope", 7.8, 5.1, 4.1, 0.25, size=13, color=NAVY, bold=True)
    _text_box(slide, "Intention and action-effect have targeted qualitative ablations. Habit is representationally useful, but still needs a routine diary workload.", 7.8, 5.42, 4.1, 0.45, size=10.5, color=INK)
    return slide


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Recall: four channels become one ranked evidence list", "method")
    channels = [
        ("Semantic", "similar meaning", BLUE),
        ("BM25", "lexical match", TEAL),
        ("Graph", "linked evidence", GOLD),
        ("Temporal", "time windows", GREEN),
    ]
    for i, (name, desc, color) in enumerate(channels):
        y = 1.35 + i * 1.08
        _card(slide, 0.9, y, 2.6, 0.76, fill=RGBColor(255, 255, 255), line=color)
        _text_box(slide, name, 1.1, y + 0.13, 2.15, 0.22, size=13.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, desc, 1.1, y + 0.41, 2.15, 0.2, size=9.2, color=MUTED, align=PP_ALIGN.CENTER)
        _arrow(slide, 3.55, y + 0.38, 5.4, 3.2, color=color, width=1.3)
    _card(slide, 5.35, 2.2, 2.5, 2.05, fill=RGBColor(255, 255, 255), line=BLUE)
    _text_box(slide, "Adaptive routing", 5.6, 2.45, 2.0, 0.28, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, "query type changes channel weights", 5.62, 2.93, 1.96, 0.48, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _text_box(slide, "semantic | temporal | causal | prospective | preference | multi-hop", 5.62, 3.5, 1.96, 0.42,
              size=8.5, color=TEAL, align=PP_ALIGN.CENTER)
    _arrow(slide, 7.86, 3.2, 9.05, 3.2, color=MUTED, width=1.8)
    _card(slide, 9.1, 1.45, 2.95, 3.75, fill=RGBColor(255, 255, 255), line=LINE)
    _text_box(slide, "Top evidence", 9.42, 1.72, 2.3, 0.3, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    for i, (txt, color) in enumerate([("fact + snippet", TEAL), ("dated memory", GREEN), ("linked node", GOLD), ("lexical hit", BLUE)]):
        _label(slide, f"{i+1}. {txt}", 9.55, 2.25 + i * 0.58, 2.05, 0.34, fill=color, size=9.2)
    _text_box(slide, "Key claim: graph is important, but the final system is deliberately multi-channel.", 1.0, 6.3, 10.95, 0.42,
              size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Graph recall: SUM instead of MAX", "method")
    _card(slide, 0.78, 1.35, 5.45, 4.75)
    _card(slide, 6.65, 1.35, 5.45, 4.75)
    _label(slide, "MAX baseline", 1.0, 1.58, 1.55, 0.42, fill=RED)
    _label(slide, "SUM activation", 6.88, 1.58, 1.75, 0.42, fill=TEAL)
    _text_box(slide, "Keeps the strongest single path", 1.0, 2.14, 4.95, 0.3, size=14, color=NAVY, bold=True)
    _text_box(slide, "Accumulates many weak paths converging on the same answer node", 6.88, 2.14, 4.85, 0.3,
              size=14, color=NAVY, bold=True)
    # MAX diagram
    coords = [(1.55, 4.6), (2.35, 3.25), (3.2, 4.65), (4.55, 3.65), (5.1, 4.85)]
    for i, (x, y) in enumerate(coords):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _x(x), _y(y), _w(0.35), _h(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RED if i == 3 else RGBColor(215, 222, 230)
        dot.line.color.rgb = dot.fill.fore_color.rgb
    _arrow(slide, 1.9, 4.78, 4.55, 3.82, color=RED, width=2.0)
    _text_box(slide, "A(v) = max(signal)", 1.25, 5.35, 4.6, 0.3, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    # SUM diagram
    answer = slide.shapes.add_shape(MSO_SHAPE.OVAL, _x(9.2), _y(3.7), _w(0.62), _h(0.62))
    answer.fill.solid()
    answer.fill.fore_color.rgb = TEAL
    answer.line.color.rgb = TEAL
    for x, y, c in [(7.45, 3.05, BLUE), (7.7, 4.72, GREEN), (10.65, 3.15, GOLD), (10.55, 4.85, ORANGE)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _x(x), _y(y), _w(0.36), _h(0.36))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.color.rgb = c
        _arrow(slide, x + 0.36, y + 0.18, 9.2, 4.01, color=c, width=1.7)
    _text_box(slide, "A(v) += total incoming signal", 7.05, 5.35, 4.6, 0.3, size=17, color=TEAL, bold=True,
              align=PP_ALIGN.CENTER)
    _label(slide, "cycle guard 1: refractory", 1.05, 6.32, 2.4, 0.34, fill=BLUE, size=8.8)
    _label(slide, "cycle guard 2: firing quota", 4.08, 6.32, 2.4, 0.34, fill=GOLD, size=8.8)
    _label(slide, "cycle guard 3: saturation", 7.08, 6.32, 2.4, 0.34, fill=GREEN, size=8.8)
    return slide


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Running example: before-travel question", "example")
    _card(slide, 0.85, 1.22, 11.25, 0.75, fill=RGBColor(255, 255, 255), line=TEAL)
    _text_box(slide, "Query: Which city was the user in before traveling to Chicago?", 1.08, 1.42, 10.8, 0.3,
              size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # timeline
    _arrow(slide, 1.15, 3.18, 11.5, 3.18, color=MUTED, width=1.6)
    events = [
        ("Memory A", "Seattle", "earlier dated evidence", TEAL),
        ("Memory B", "Chicago", "anchor travel evidence", GOLD),
        ("Answer", "Seattle", "before Chicago", GREEN),
    ]
    xs = [2.15, 6.0, 9.85]
    for x, (label, city, note, color) in zip(xs, events):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _x(x), _y(3.0), _w(0.36), _h(0.36))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        _card(slide, x - 0.7, 3.62, 1.75, 1.05, fill=RGBColor(255, 255, 255), line=color)
        _text_box(slide, label, x - 0.5, 3.78, 1.35, 0.22, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, city, x - 0.5, 4.08, 1.35, 0.24, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _text_box(slide, note, x - 0.5, 4.38, 1.35, 0.2, size=8.2, color=MUTED, align=PP_ALIGN.CENTER)
    _card(slide, 1.0, 5.38, 3.1, 0.88)
    _card(slide, 4.9, 5.38, 3.1, 0.88)
    _card(slide, 8.8, 5.38, 3.1, 0.88)
    _text_box(slide, "1. Detect relation", 1.22, 5.52, 2.65, 0.22, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, "before + travel + city", 1.22, 5.82, 2.65, 0.22, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    _text_box(slide, "2. Recall dated memories", 5.12, 5.52, 2.65, 0.22, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, "Seattle and Chicago evidence", 5.12, 5.82, 2.65, 0.22, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    _text_box(slide, "3. Guard generation", 9.02, 5.52, 2.65, 0.22, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, "answer from evidence only", 9.02, 5.82, 2.65, 0.22, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    return slide


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Evaluation: manual verdicts, not blind judge trust", "evaluation")
    _fit_picture(slide, IMG_DIR / "manual_evaluation_flow.png", 0.6, 1.08, 7.0, 4.95)
    _card(slide, 7.9, 1.35, 3.95, 4.65)
    _text_box(slide, "Final scoring rules", 8.18, 1.62, 3.4, 0.32, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _bullet_box(slide, [
        "Automatic judge is diagnostic only.",
        "Manual PASS is the main metric.",
        "LoCoMo PASS means fully correct or core-correct and usable.",
        "Blank-gold cases pass only when memory insufficiency is explicit.",
    ], 8.25, 2.25, 3.2, 2.45, size=12.5)
    _label(slide, "Question", 0.98, 6.04, 1.2, 0.34, fill=BLUE, size=8.8)
    _label(slide, "CogMem answer", 2.45, 6.04, 1.55, 0.34, fill=TEAL, size=8.8)
    _label(slide, "Human verdict", 4.25, 6.04, 1.5, 0.34, fill=GOLD, size=8.8)
    _label(slide, "Category totals", 6.0, 6.04, 1.55, 0.34, fill=GREEN, size=8.8)
    return slide


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "LongMemEval v16: multi-channel recall is strongest", "results")
    _stat_card(slide, "Best multi-channel", "29/35", "manual PASS = 82.9%", 0.85, 1.26, 3.2, 1.55, TEAL)
    _stat_card(slide, "Full six-type baseline", "26/35", "manual PASS = 74.3%", 4.75, 1.26, 3.2, 1.55, BLUE)
    _stat_card(slide, "Verified subset", "35", "LongMemEval v16 questions", 8.65, 1.26, 3.2, 1.55, GOLD)
    _card(slide, 1.15, 3.35, 10.4, 2.45)
    _bar(slide, "Best multi-channel", 82.9, 1.55, 3.72, 5.75, TEAL)
    _bar(slide, "Full six-type baseline", 74.3, 1.55, 4.32, 5.75, BLUE)
    _text_box(slide, "Graph-only controls", 1.55, 4.92, 2.4, 0.28, size=12, color=INK, bold=True)
    _text_box(slide, "lower-band controls used for isolation, not the final system", 4.05, 4.92, 5.65, 0.28,
              size=10.3, color=MUTED)
    _text_box(slide, "Takeaway: CogMem works best when graph structure complements lexical, semantic, and temporal evidence.", 1.05, 6.35, 10.9, 0.35,
              size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Control experiment: SUM improves graph-only top-k priority", "results")
    _stat_card(slide, "Graph-only SUM recall@5", "0.8052", "mean session recall@5", 0.78, 1.22, 3.2, 1.55, TEAL)
    _stat_card(slide, "Graph-only MAX recall@5", "0.7624", "mean session recall@5", 4.3, 1.22, 3.2, 1.55, RED)
    _stat_card(slide, "Absolute lift", "+0.0429", "SUM minus MAX at top-5", 7.82, 1.22, 3.2, 1.55, GOLD)
    _card(slide, 0.92, 3.22, 5.2, 2.2)
    _card(slide, 6.72, 3.22, 5.2, 2.2)
    _text_box(slide, "Rank budget view", 1.2, 3.48, 4.6, 0.26, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    for i in range(1, 8):
        fill = GREEN if i <= 5 else RGBColor(225, 232, 240)
        _label(slide, str(i), 1.15 + (i - 1) * 0.62, 4.05, 0.38, 0.38, fill=fill, size=10)
    _text_box(slide, "Relevant evidence should land inside ranks 1-5.", 1.2, 4.76, 4.6, 0.28, size=11.5, color=MUTED,
              align=PP_ALIGN.CENTER)
    _text_box(slide, "Result interpretation", 7.0, 3.48, 4.6, 0.26, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _bullet_box(slide, [
        "Recall@10 is equal at 0.8481.",
        "SUM better in 2/35 questions.",
        "MAX better in 0/35 questions.",
        "Top-10 order changes in 12/35 questions.",
    ], 7.08, 3.95, 4.35, 1.2, size=11.5)
    _text_box(slide, "Honest claim: SUM does not create new evidence; it moves convergent evidence earlier in the prompt budget.",
              1.0, 6.18, 11.0, 0.45, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "LoCoMo: final system crosses the 70% target", "results")
    _stat_card(slide, "Final evidence-guard configuration", "119/161", "manual PASS = 73.9%", 0.78, 1.16, 3.5, 1.6, TEAL)
    _stat_card(slide, "Previous baseline", "97/161", "manual PASS = 60.2%", 4.8, 1.16, 3.2, 1.6, BLUE)
    _stat_card(slide, "Recovered cases", "+22", "additional manual PASS", 8.52, 1.16, 3.1, 1.6, GOLD)
    _card(slide, 1.0, 3.15, 10.9, 2.9)
    categories = [
        ("multi-hop", 91.7, TEAL),
        ("causal", 90.9, GREEN),
        ("preference", 88.2, GOLD),
        ("single-hop", 71.6, BLUE),
        ("temporal", 41.7, RED),
    ]
    for i, (label, val, color) in enumerate(categories):
        _bar(slide, label, val, 1.38, 3.55 + i * 0.45, 6.5, color)
    _text_box(slide, "Takeaway: the full system is useful on a harder long-dialogue benchmark, but temporal reasoning still needs focused work.",
              1.0, 6.42, 11.0, 0.32, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Qualitative proof: intention stores unfinished plans", "typed memory proof")
    _card(slide, 0.78, 1.18, 11.25, 0.72, fill=RGBColor(255, 255, 255), line=GOLD)
    _text_box(slide, "Target question: What sustainability habit did the user intend to start but has not?",
              1.05, 1.38, 10.75, 0.28, size=16.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    _card(slide, 0.78, 2.15, 5.35, 3.1, fill=RGBColor(247, 253, 250), line=GREEN)
    _label(slide, "Full memory bank", 1.05, 2.38, 2.1, 0.36, fill=GREEN, size=9.5)
    _text_box(slide, "37 facts, including 4 intention nodes", 3.32, 2.42, 2.45, 0.25,
              size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _mini_node(slide, "plans\ncomposting", 1.15, 3.15, RED, w=1.25, h=0.55)
    _mini_node(slide, "has not\nstarted", 2.55, 3.15, RED, w=1.22, h=0.55)
    _mini_node(slide, "trash bags", 1.85, 4.08, TEAL, w=1.18, h=0.42)
    _arrow(slide, 2.4, 3.42, 2.55, 3.42, color=RED, width=1.4)
    _arrow(slide, 2.42, 4.12, 3.62, 4.12, color=MUTED, width=1.1)
    _label(slide, "Answer: Composting", 3.7, 3.78, 1.85, 0.48, fill=GOLD, size=9.6)
    _text_box(slide, "The gold plan remains physically present as intention-typed evidence.",
              1.1, 4.67, 4.65, 0.28, size=10.5, color=INK, align=PP_ALIGN.CENTER)

    _card(slide, 6.7, 2.15, 5.35, 3.1, fill=RGBColor(255, 248, 248), line=RED)
    _label(slide, "Intention-ablated bank", 6.98, 2.38, 2.45, 0.36, fill=RED, size=9.5)
    _text_box(slide, "26 facts, 0 intention nodes", 9.62, 2.42, 1.95, 0.25,
              size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _mini_node(slide, "rain barrel", 7.15, 3.2, BLUE, w=1.18, h=0.42)
    _mini_node(slide, "garden", 8.62, 3.2, BLUE, w=1.05, h=0.42)
    _mini_node(slide, "no red\nplan nodes", 10.08, 3.05, RED, w=1.28, h=0.58)
    _arrow(slide, 8.35, 3.41, 8.62, 3.41, color=BLUE, width=1.2)
    _arrow(slide, 9.68, 3.41, 10.08, 3.34, color=MUTED, width=1.0)
    _label(slide, "Wrong decoy: rainwater", 8.08, 4.02, 2.15, 0.46, fill=RED, size=9.0)
    _text_box(slide, "The system still recalls evidence, but the missing type forces a topic swap.",
              7.02, 4.67, 4.65, 0.28, size=10.5, color=INK, align=PP_ALIGN.CENTER)

    _card(slide, 0.95, 5.72, 10.95, 0.84, fill=RGBColor(255, 252, 244), line=GOLD)
    _text_box(slide, "Interpretation: intention is not universally necessary, but it is necessary in sparse plan-not-done cases where no experience fact carries the same content.",
              1.2, 5.94, 10.45, 0.32, size=12.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "Qualitative proof: action-effect stores tool outcomes", "typed memory proof")
    _card(slide, 0.78, 1.14, 11.25, 0.78, fill=RGBColor(255, 255, 255), line=PURPLE)
    _text_box(slide, "Target question: When Stripe returns HTTP 429 with Retry-After, what does the agent do and what happens?",
              1.05, 1.34, 10.75, 0.32, size=15.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    _card(slide, 0.78, 2.14, 5.35, 3.18, fill=RGBColor(249, 247, 253), line=PURPLE)
    _label(slide, "Full memory bank", 1.05, 2.38, 2.1, 0.36, fill=PURPLE, size=9.5)
    _text_box(slide, "13 facts, including 7 action-effect nodes", 3.3, 2.42, 2.45, 0.25,
              size=10.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _mini_node(slide, "HTTP 429\nRetry-After", 1.08, 3.12, BLUE, w=1.35, h=0.58)
    _mini_node(slide, "sleep +\nretry", 2.8, 3.12, PURPLE, w=1.18, h=0.58)
    _mini_node(slide, "returns\n200", 4.35, 3.12, GREEN, w=1.05, h=0.58)
    _arrow(slide, 2.43, 3.41, 2.8, 3.41, color=PURPLE, width=1.5)
    _arrow(slide, 3.98, 3.41, 4.35, 3.41, color=PURPLE, width=1.5)
    _label(slide, "Answer: backoff respects Retry-After; later calls succeed", 1.12, 4.18, 4.75, 0.5,
           fill=GOLD, size=8.4)
    _text_box(slide, "The causal triple is available: precondition -> action -> outcome.",
              1.12, 4.84, 4.65, 0.25, size=10.5, color=INK, align=PP_ALIGN.CENTER)

    _card(slide, 6.7, 2.14, 5.35, 3.18, fill=RGBColor(255, 248, 248), line=RED)
    _label(slide, "Action-effect-ablated bank", 6.96, 2.38, 2.88, 0.36, fill=RED, size=9.0)
    _text_box(slide, "12 facts, 0 action-effect nodes", 10.04, 2.42, 1.55, 0.25,
              size=10.1, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _mini_node(slide, "429 has\nheader", 7.05, 3.1, BLUE, w=1.18, h=0.58)
    _mini_node(slide, "backoff\nadded", 8.62, 3.1, TEAL, w=1.12, h=0.58)
    _mini_node(slide, "specific\nsleep rule?", 10.05, 3.1, RED, w=1.22, h=0.58)
    _arrow(slide, 8.23, 3.39, 8.62, 3.39, color=MUTED, width=1.0)
    _arrow(slide, 9.74, 3.39, 10.05, 3.39, color=MUTED, width=1.0)
    _label(slide, "Answer hedges: specifics not documented", 7.35, 4.18, 3.55, 0.5,
           fill=RED, size=8.8)
    _text_box(slide, "Recall is not empty; the missing type removes the exact action-result rule.",
              7.02, 4.84, 4.65, 0.25, size=10.5, color=INK, align=PP_ALIGN.CENTER)

    _card(slide, 0.95, 5.74, 10.95, 0.84, fill=RGBColor(255, 252, 244), line=GOLD)
    _text_box(slide, "Interpretation: action-effect cleanly discriminates in 5/12 mocked agentic traces. This is conditional evidence, not a universal claim.",
              1.2, 5.96, 10.45, 0.32, size=12.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bg(slide, prs, "What was proved, and what remains", "contribution proof")
    rows = [
        ("Multi-channel memory recall", "LongMemEval 29/35; LoCoMo 119/161", "Strong"),
        ("SUM graph activation", "0.8052 vs 0.7624 recall@5", "Strong, scoped"),
        ("Evidence guard + snippets", "+22 LoCoMo PASS cases", "Strong practical"),
        ("Intention / action-effect types", "paired qualitative ablations: composting and HTTP 429", "Conditional"),
        ("Habit Network", "needs routine diary workload", "Limitation"),
    ]
    x0, y0 = 0.75, 1.25
    widths = [3.7, 5.0, 2.15]
    headers = ["Claim", "Evidence", "Status"]
    for i, header in enumerate(headers):
        _label(slide, header, x0 + sum(widths[:i]), y0, widths[i] - 0.05, 0.42, fill=NAVY, size=10.5)
    for r, row in enumerate(rows):
        y = y0 + 0.52 + r * 0.72
        status_color = GREEN if "Strong" in row[2] else GOLD if "Conditional" in row[2] else RED
        for c, txt in enumerate(row):
            _card(slide, x0 + sum(widths[:c]), y, widths[c] - 0.05, 0.58, fill=RGBColor(255, 255, 255),
                  line=LINE)
            _text_box(slide, txt, x0 + sum(widths[:c]) + 0.08, y + 0.11, widths[c] - 0.2, 0.28,
                      size=10.4, color=status_color if c == 2 else INK, bold=c in (0, 2),
                      align=PP_ALIGN.CENTER if c == 2 else PP_ALIGN.LEFT)
    _card(slide, 0.95, 5.72, 10.95, 0.86, fill=RGBColor(255, 252, 244), line=GOLD)
    _text_box(slide, "Future work should move toward realistic assistant scenarios: plan lifecycle tasks, dense tool-action traces, and multi-week habit diaries.",
              1.2, 5.95, 10.45, 0.32, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _text_box(slide, "THANK\nYOU !", 6.0, 2.45, 3.0, 1.25, size=38, color=RGBColor(196, 0, 0), bold=True,
              align=PP_ALIGN.CENTER, font="Aptos Display")
    _text_box(slide, "Questions?", 6.05, 4.0, 2.9, 0.34, size=18, color=INK, bold=True,
              align=PP_ALIGN.CENTER)
    _text_box(slide, "Code, experiments, manual verdicts, and thesis artifacts are preserved in the repository.",
              5.35, 4.85, 4.2, 0.44, size=10.8, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    return slide


def rebuild() -> None:
    prs = Presentation(TEMPLATE)
    _clear_slides(prs)
    for builder in [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
        slide_15,
        slide_16,
    ]:
        builder(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Rebuilt {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    rebuild()
